"""PowerPoint presentation (.pptx) to Markdown converter."""

from io import BytesIO

from pptx import Presentation
from pptx.shapes.base import BaseShape
from pptx.table import Table

from .base import ConversionError, ConversionResult, strip_extension


def convert_pptx(file_bytes: bytes, filename: str) -> ConversionResult:
    """Convert a PowerPoint presentation to Markdown.

    Extracts:
    - Slide titles as headings
    - Text from all shapes
    - Tables
    - Speaker notes

    Args:
        file_bytes: Raw .pptx file content
        filename: Original filename

    Returns:
        ConversionResult with Markdown content
    """
    try:
        prs = Presentation(BytesIO(file_bytes))
    except Exception as e:
        raise ConversionError(f"Failed to parse PowerPoint file: {e}") from e

    markdown_sections: list[str] = []

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_md = _slide_to_markdown(slide, slide_num)
        if slide_md:
            markdown_sections.append(slide_md)

    if not markdown_sections:
        markdown_content = f"# {strip_extension(filename)}\n\n*Empty presentation*"
    else:
        markdown_content = "\n\n---\n\n".join(markdown_sections)

    new_filename = f"{strip_extension(filename)}.md"

    return ConversionResult(
        content=markdown_content.encode("utf-8"),
        mime_type="text/markdown",
        filename=new_filename,
    )


def _slide_to_markdown(slide, slide_num: int) -> str | None:
    """Convert a slide to Markdown."""
    lines: list[str] = []

    # Try to get slide title
    title = None
    if slide.shapes.title:
        title = slide.shapes.title.text.strip()

    if title:
        lines.append(f"## Slide {slide_num}: {title}")
    else:
        lines.append(f"## Slide {slide_num}")

    lines.append("")

    # Extract text from all shapes
    for shape in slide.shapes:
        shape_md = _shape_to_markdown(shape)
        if shape_md:
            lines.append(shape_md)

    # Extract speaker notes
    if slide.has_notes_slide:
        notes_frame = slide.notes_slide.notes_text_frame
        if notes_frame and notes_frame.text.strip():
            lines.append("")
            lines.append("**Speaker Notes:**")
            lines.append(f"> {notes_frame.text.strip()}")

    # Check if slide has any content beyond the heading
    content_lines = [line for line in lines[2:] if line.strip()]
    if not content_lines:
        return None

    return "\n".join(lines)


def _shape_to_markdown(shape: BaseShape) -> str | None:
    """Extract text from a shape."""
    # Handle tables
    if shape.has_table:
        return _table_to_markdown(shape.table)

    # Handle text frames
    if shape.has_text_frame:
        paragraphs: list[str] = []
        for para in shape.text_frame.paragraphs:
            text = para.text.strip()
            if text:
                # Check indentation level for lists
                if para.level > 0:
                    indent = "  " * para.level
                    paragraphs.append(f"{indent}- {text}")
                else:
                    paragraphs.append(text)

        if paragraphs:
            return "\n".join(paragraphs)

    return None


def _table_to_markdown(table: Table) -> str:
    """Convert a PowerPoint table to Markdown."""
    rows: list[list[str]] = []

    for row in table.rows:
        cells = []
        for cell in row.cells:
            text = cell.text.strip().replace("\n", " ").replace("|", "\\|")
            cells.append(text)
        rows.append(cells)

    if not rows:
        return ""

    lines: list[str] = []

    # Header row
    header = rows[0]
    lines.append("| " + " | ".join(header) + " |")

    # Separator
    lines.append("| " + " | ".join(["---"] * len(header)) + " |")

    # Data rows
    for row in rows[1:]:
        while len(row) < len(header):
            row.append("")
        lines.append("| " + " | ".join(row[: len(header)]) + " |")

    return "\n".join(lines)
